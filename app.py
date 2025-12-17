import streamlit as st
from openai import OpenAI
import PyPDF2
from docx import Document
from pptx import Presentation
import pandas as pd
from PIL import Image
import pytesseract
import speech_recognition as sr
from pydub import AudioSegment
import io
import os

# --- 1. 頁面設定 ---
st.set_page_config(page_title="DSE 全能温習平台", layout="wide", page_icon="🇭🇰")

# --- 2. 安全讀取 API Key ---
api_key = None
if "DEEPSEEK_API_KEY" in st.secrets:
    api_key = st.secrets["DEEPSEEK_API_KEY"]
else:
    api_key = st.sidebar.text_input("DeepSeek API Key", type="password")

if not api_key:
    st.warning("⚠️ 請設定 API Key 以繼續。")
    st.stop()

client = OpenAI(api_key=api_key, base_url="https://api.deepseek.com")

# --- 3. 核心功能：萬能檔案讀取器 ---
def extract_text_from_file(uploaded_file):
    text = ""
    file_type = uploaded_file.name.split('.')[-1].lower()
    
    try:
        # A. 處理 PDF
        if file_type == 'pdf':
            reader = PyPDF2.PdfReader(uploaded_file)
            for page in reader.pages:
                text += page.extract_text() + "\n"
                
        # B. 處理 Word
        elif file_type == 'docx':
            doc = Document(uploaded_file)
            for para in doc.paragraphs:
                text += para.text + "\n"
                
        # C. 處理 PowerPoint
        elif file_type == 'pptx':
            prs = Presentation(uploaded_file)
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text += shape.text + "\n"
        
        # D. 處理 Excel / CSV (轉換為 Markdown 表格)
        elif file_type in ['xlsx', 'xls', 'csv']:
            if file_type == 'csv':
                df = pd.read_csv(uploaded_file)
            else:
                df = pd.read_excel(uploaded_file)
            # 將表格轉為文字格式，讓 AI 讀懂
            text = df.to_markdown(index=False)

        # E. 處理 圖片 (OCR)
        elif file_type in ['jpg', 'jpeg', 'png']:
            image = Image.open(uploaded_file)
            # 使用 Tesseract 進行 OCR，設定繁體中文+英文
            # 注意：Streamlit Cloud 需透過 packages.txt 安裝 tesseract
            text = pytesseract.image_to_string(image, lang='chi_tra+eng')
            if not text.strip():
                text = "[OCR 提示] 圖片中未能識別出文字，請確保圖片清晰。"

        # F. 處理 聲音 (Speech to Text)
        elif file_type in ['mp3', 'wav', 'm4a']:
            # 轉換音訊格式為 wav (SpeechRecognition 需要 wav)
            audio = AudioSegment.from_file(uploaded_file)
            wav_buffer = io.BytesIO()
            audio.export(wav_buffer, format="wav")
            wav_buffer.seek(0)
            
            recognizer = sr.Recognizer()
            with sr.AudioFile(wav_buffer) as source:
                audio_data = recognizer.record(source)
                # 使用 Google 免費 API 識別 (支援廣東話/繁中)
                # 廣東話 code: yue-Hant-HK, 國語: cmn-Hant-TW
                try:
                    text = recognizer.recognize_google(audio_data, language="yue-Hant-HK")
                except sr.UnknownValueError:
                    text = "[Audio 提示] 無法識別語音，可能是聲音太小或雜訊過多。"
                except sr.RequestError:
                    text = "[Audio 提示] 語音識別服務暫時無法連接。"

        # G. 處理 純文字
        elif file_type in ['txt', 'md']:
            text = uploaded_file.read().decode("utf-8")
            
    except Exception as e:
        return f"讀取檔案錯誤 ({file_type}): {str(e)}"
        
    return text

# --- 4. 側邊欄 ---
with st.sidebar:
    st.title("🇭🇰 DSE 備戰中心")
    st.caption("支援: PDF, Word, PPT, Excel, 圖片, 聲音")
    st.divider()
    subject = st.selectbox("科目", ["Biology", "Chemistry", "Economics", "Chinese", "English", "History", "Maths"])
    
# --- 5. 主功能區 ---
tab_factory, tab_study = st.tabs(["🛠️ 步驟一：萬能資料清洗", "🎓 步驟二：智能溫習室"])

# ==========================================
# TAB 1: 資料清洗 (支援 OCR, STT, Pandas)
# ==========================================
with tab_factory:
    st.header(f"🧹 {subject} - 多媒體資料處理工廠")
    st.info("支援檔案：PDF, Word, PPT, Excel, CSV, 圖片 (JPG/PNG), 錄音 (MP3/WAV)")
    
    uploaded_file = st.file_uploader(
        "上傳任何教材檔案", 
        type=["pdf", "docx", "pptx", "xlsx", "csv", "txt", "md", "jpg", "png", "mp3", "wav", "m4a"]
    )
    
    if uploaded_file:
        with st.spinner("正在讀取檔案內容 (圖片/聲音可能需要較長時間)..."):
            raw_text = extract_text_from_file(uploaded_file)
        
        # 檢查是否讀取成功
        if str(raw_text).startswith("讀取檔案錯誤") or not raw_text:
            st.error(f"無法提取內容：{raw_text}")
        else:
            word_count = len(str(raw_text))
            st.success(f"✅ 成功提取內容！共 {word_count} 字。")
            with st.expander("預覽提取的原始文字"):
                st.text(raw_text[:2000] + "..." if word_count > 2000 else raw_text)
            
            if st.button("🚀 交給 DeepSeek 整理重點"):
                with st.spinner("DeepSeek 正在分析並整理筆記..."):
                    # 根據檔案類型微調 Prompt
                    clean_prompt = f"""
                    你是一位專業的 DSE {subject} 導師。
                    請處理以下原始文本 (可能來自圖片 OCR、錄音轉錄或 Excel 表格)，整理為高品質的 Markdown 筆記。
                    
                    原始文本：
                    {raw_text[:15000]} 
                    
                    要求：
                    1. 【修正錯誤】：如果是 OCR 或語音轉錄，請自動修正明顯的錯字或語意不通處。
                    2. 【結構化】：按課題分類，如果是表格數據，請整理為分析結論。
                    3. 【精煉】：保留 DSE 考試重點 (Keywords)。
                    """
                    
                    try:
                        response = client.chat.completions.create(
                            model="deepseek-chat",
                            messages=[{"role": "user", "content": clean_prompt}]
                        )
                        cleaned_content = response.choices[0].message.content
                        
                        st.subheader("📝 整理後的筆記")
                        st.text_area("Result", cleaned_content, height=400)
                        
                        file_name = f"{subject}_Cleaned_Notes.txt"
                        st.download_button(
                            label="📥 下載筆記 (.txt) -> 用於 NotebookLM",
                            data=cleaned_content,
                            file_name=file_name,
                            mime="text/plain"
                        )
                    except Exception as e:
                        st.error(f"API 呼叫錯誤: {e}")

# ==========================================
# TAB 2: 溫習室 (邏輯保持不變)
# ==========================================
with tab_study:
    st.header(f"🎓 {subject} - 衝刺模式")
    
    col_input, col_main = st.columns([1, 2])
    with col_input:
        st.markdown("### 載入資源")
        notes_file = st.file_uploader("上傳清洗後的筆記 (.txt)", type=["txt", "md", "docx"], key="notes")
        audio_file = st.file_uploader("上傳 NotebookLM 音檔 (.mp3)", type=["mp3", "wav"], key="audio")
        
        notes_text = ""
        if notes_file:
             notes_text = extract_text_from_file(notes_file)
    
    with col_main:
        if not notes_text:
            st.info("👈 請先上傳筆記")
        else:
            sub_tab1, sub_tab2 = st.tabs(["💬 導師問答", "✍️ 模擬試卷"])
            
            with sub_tab1:
                st.subheader("AI 導師")
                if "messages" not in st.session_state:
                    st.session_state.messages = []
                for msg in st.session_state.messages:
                    st.chat_message(msg["role"]).write(msg["content"])
                
                if user_input := st.chat_input("輸入問題..."):
                    st.session_state.messages.append({"role": "user", "content": user_input})
                    st.chat_message("user").write(user_input)
                    with st.chat_message("assistant"):
                        rag_prompt = f"你是 DSE 導師。根據筆記回答 (廣東話)：\n{notes_text[:10000]}"
                        stream = client.chat.completions.create(
                            model="deepseek-chat",
                            messages=[
                                {"role": "system", "content": rag_prompt},
                                {"role": "user", "content": user_input}
                            ],
                            stream=True
                        )
                        response = st.write_stream(stream)
                    st.session_state.messages.append({"role": "assistant", "content": response})

            with sub_tab2:
                if st.button("生成題目"):
                     with st.spinner("出卷中..."):
                        res = client.chat.completions.create(
                            model="deepseek-chat",
                            messages=[{"role": "user", "content": f"根據筆記出 DSE 題目：{notes_text[:5000]}"}]
                        )
                        st.markdown(res.choices[0].message.content)                    """
                    
                    response = client.chat.completions.create(
                        model="deepseek-chat",
                        messages=[{"role": "user", "content": clean_prompt}]
                    )
                    cleaned_content = response.choices[0].message.content
                    
                    # 顯示結果與下載
                    st.subheader("📝 筆記預覽")
                    st.text_area("Result", cleaned_content, height=300)
                    
                    file_name = f"{subject}_Cleaned_Notes.txt"
                    st.download_button(
                        label="📥 下載 .txt 檔案 (用於 NotebookLM)",
                        data=cleaned_content,
                        file_name=file_name,
                        mime="text/plain"
                    )
                    st.success("✅ 完成！請將此檔案上傳至 NotebookLM 生成 Audio。")
                    
        except Exception as e:
            st.error(f"讀取 PDF 時發生錯誤: {e}")

# ==========================================
# TAB 2: 智能溫習室 (Study Room)
# ==========================================
with tab_study:
    st.header(f"🎓 {subject} - 衝刺模式")
    
    col_input, col_main = st.columns([1, 2])
    
    with col_input:
        st.markdown("### 1. 載入資源")
        notes_file = st.file_uploader("上傳清洗後的筆記 (.txt)", type=["txt", "md"], key="notes")
        audio_file = st.file_uploader("上傳 NotebookLM 音檔 (.mp3)", type=["mp3", "wav"], key="audio")
        
        notes_text = ""
        if notes_file:
            notes_text = notes_file.read().decode("utf-8")
    
    with col_main:
        if not notes_text:
            st.info("👈 請先在左側上傳資源 (步驟一產出的 TXT + NotebookLM 的 MP3)")
        else:
            # === 功能分頁 ===
            sub_tab1, sub_tab2, sub_tab3 = st.tabs(["🎧 聽覺學習", "💬 導師問答", "✍️ 模擬試卷"])
            
            # --- 聽覺學習 ---
            with sub_tab1:
                st.subheader("NotebookLM Podcast")
                if audio_file:
                    st.audio(audio_file)
                else:
                    st.warning("未上傳音頻，建議配合 NotebookLM 使用以達最佳效果。")
                
                with st.expander("查看完整筆記內容"):
                    st.markdown(notes_text)

            # --- 導師問答 ---
            with sub_tab2:
                st.subheader("AI 導師 (DeepSeek)")
                if "messages" not in st.session_state:
                    st.session_state.messages = []

                for msg in st.session_state.messages:
                    st.chat_message(msg["role"]).write(msg["content"])

                if user_input := st.chat_input("e.g. 用廣東話解釋呢個 Concept"):
                    st.session_state.messages.append({"role": "user", "content": user_input})
                    st.chat_message("user").write(user_input)

                    with st.chat_message("assistant"):
                        rag_prompt = f"""
                        你是一位 DSE {subject} 導師。請根據以下筆記回答學生問題。
                        必須使用【廣東話】。
                        筆記內容：{notes_text[:12000]}
                        """
                        stream = client.chat.completions.create(
                            model="deepseek-chat",
                            messages=[
                                {"role": "system", "content": rag_prompt},
                                {"role": "user", "content": user_input}
                            ],
                            stream=True
                        )
                        response = st.write_stream(stream)
                    st.session_state.messages.append({"role": "assistant", "content": response})

            # --- 模擬試卷 ---
            with sub_tab3:
                st.subheader("DSE 題目生成器")
                diff = st.select_slider("難度", options=["Level 2", "Level 4", "Level 5**"])
                if st.button("生成題目"):
                    with st.spinner("出卷中..."):
                        q_prompt = f"""
                        根據筆記，設計一條 {subject} {diff} 的題目 (MC 或 LQ)。
                        附帶詳細 Marking Scheme。
                        筆記：{notes_text[:5000]}
                        """
                        res = client.chat.completions.create(
                            model="deepseek-chat",
                            messages=[{"role": "user", "content": q_prompt}]
                        )
                        st.markdown(res.choices[0].message.content)
